import uuid
from io import BytesIO
from unittest.mock import AsyncMock, MagicMock

import pytest
from pptx import Presentation

from app.models.document import Document
from app.models.document_version import DocumentVersion
from app.models.ingestion_job import IngestionJob
from app.services.document_processing_exceptions import PptxTextExtractionError
from app.services.processing_status import ProcessingStage
from app.services.processors.base import ProcessingContext
from app.services.processors.pptx_text_extraction import PptxTextExtractionProcessor


def _build_pptx(text: str) -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = text
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


@pytest.fixture
def pptx_bytes() -> bytes:
    return _build_pptx("Processor test content")


@pytest.fixture
def processing_context() -> ProcessingContext:
    document_id = uuid.uuid4()
    version_id = uuid.uuid4()
    document = Document(
        id=document_id,
        title="Presentation",
        original_filename="training.pptx",
        doc_type="manual",
        status="processing",
        uploaded_by=None,
        extra_metadata={},
    )
    version = DocumentVersion(
        id=version_id,
        document_id=document_id,
        version_no=1,
        storage_uri="documents/test/v1/training.pptx",
        checksum_sha256="abc123",
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        file_extension="pptx",
        file_size_bytes=1024,
        is_latest=True,
    )
    job = IngestionJob(
        id=uuid.uuid4(),
        document_id=document_id,
        status="processing",
        stage=ProcessingStage.PROCESSING.value,
    )
    return ProcessingContext(document=document, version=version, job=job)


@pytest.fixture
def mock_storage(pptx_bytes: bytes) -> MagicMock:
    storage = MagicMock()
    storage.read.return_value = pptx_bytes
    return storage


@pytest.fixture
def mock_repository() -> AsyncMock:
    repository = AsyncMock()
    repository.upsert_extracted_text.return_value = MagicMock()
    repository.update_document_version.return_value = MagicMock()
    repository.update_ingestion_job.return_value = MagicMock()
    return repository


@pytest.fixture
def processor(mock_storage: MagicMock, mock_repository: AsyncMock) -> PptxTextExtractionProcessor:
    return PptxTextExtractionProcessor(
        storage=mock_storage,
        document_repository=mock_repository,
    )


@pytest.mark.asyncio
async def test_pptx_processor_extracts_and_persists_text(
    processor: PptxTextExtractionProcessor,
    mock_storage: MagicMock,
    mock_repository: AsyncMock,
    processing_context: ProcessingContext,
) -> None:
    await processor.process(processing_context)

    mock_storage.read.assert_called_once_with(processing_context.version.storage_uri)
    mock_repository.update_ingestion_job.assert_awaited_with(
        processing_context.job.id,
        stage=ProcessingStage.TEXT_EXTRACTION.value,
    )
    saved_payload = mock_repository.upsert_extracted_text.await_args.kwargs
    assert "Processor test content" in saved_payload["full_text"]
    assert saved_payload["extraction_method"] == "python-pptx"
    assert saved_payload["requires_ocr"] is False
    mock_repository.update_document_version.assert_awaited_with(
        processing_context.version.id,
        page_count=1,
    )


@pytest.mark.asyncio
async def test_pptx_processor_skips_non_pptx_documents(
    processor: PptxTextExtractionProcessor,
    mock_storage: MagicMock,
    mock_repository: AsyncMock,
    processing_context: ProcessingContext,
) -> None:
    processing_context.version.file_extension = "pdf"

    await processor.process(processing_context)

    mock_storage.read.assert_not_called()
    mock_repository.upsert_extracted_text.assert_not_awaited()


@pytest.mark.asyncio
async def test_pptx_processor_raises_when_storage_read_fails(
    processor: PptxTextExtractionProcessor,
    mock_storage: MagicMock,
    processing_context: ProcessingContext,
) -> None:
    from app.core.storage.exceptions import StorageError

    mock_storage.read.side_effect = StorageError("missing")

    with pytest.raises(PptxTextExtractionError, match="Failed to read stored PPTX"):
        await processor.process(processing_context)
