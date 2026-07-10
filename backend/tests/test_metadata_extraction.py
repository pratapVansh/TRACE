from datetime import UTC, datetime
from io import BytesIO

import fitz
import pytest
from docx import Document as DocxDocument
from openpyxl import Workbook
from pptx import Presentation

from app.services.document_processing_exceptions import MetadataExtractionError
from app.services.metadata_extraction import extract_document_metadata


def _build_pdf(author: str = "Jane Engineer") -> bytes:
    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), "Metadata test")
    document.set_metadata(
        {
            "author": author,
            "creationDate": "D:20240115103000",
            "modDate": "D:20240620154530",
        },
    )
    content = document.tobytes()
    document.close()
    return content


def _build_docx() -> bytes:
    document = DocxDocument()
    properties = document.core_properties
    properties.author = "Doc Author"
    properties.created = datetime(2024, 2, 1, 9, 0, tzinfo=UTC)
    properties.modified = datetime(2024, 2, 2, 10, 30, tzinfo=UTC)
    document.add_paragraph("Sample")
    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _build_pptx() -> bytes:
    presentation = Presentation()
    presentation.core_properties.author = "Slide Author"
    presentation.core_properties.created = datetime(2024, 3, 1, 8, 0, tzinfo=UTC)
    presentation.core_properties.modified = datetime(2024, 3, 2, 9, 15, tzinfo=UTC)
    presentation.slides.add_slide(presentation.slide_layouts[5])
    presentation.slides.add_slide(presentation.slide_layouts[5])
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _build_xlsx() -> bytes:
    workbook = Workbook()
    workbook.create_sheet("Logs")
    workbook.properties.creator = "Sheet Author"
    workbook.properties.created = datetime(2024, 4, 1, 12, 0, tzinfo=UTC)
    workbook.properties.modified = datetime(2024, 4, 2, 13, 45, tzinfo=UTC)
    buffer = BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def test_extract_document_metadata_for_pdf() -> None:
    content = _build_pdf()

    result = extract_document_metadata(
        content,
        mime_type="application/pdf",
        file_extension="pdf",
    )

    assert result.page_count == 1
    assert result.file_size_bytes == len(content)
    assert result.file_type == "application/pdf"
    assert result.file_extension == "pdf"
    assert result.author == "Jane Engineer"
    assert result.creation_date == "2024-01-15T10:30:00+00:00"
    assert result.modification_date == "2024-06-20T15:45:30+00:00"


def test_extract_document_metadata_for_docx() -> None:
    content = _build_docx()

    result = extract_document_metadata(
        content,
        mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        file_extension="docx",
    )

    assert result.author == "Doc Author"
    assert result.creation_date == "2024-02-01T09:00:00+00:00"
    assert result.modification_date == "2024-02-02T10:30:00+00:00"


def test_extract_document_metadata_for_pptx() -> None:
    content = _build_pptx()

    result = extract_document_metadata(
        content,
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        file_extension="pptx",
    )

    assert result.page_count == 2
    assert result.author == "Slide Author"


def test_extract_document_metadata_for_xlsx() -> None:
    content = _build_xlsx()

    result = extract_document_metadata(
        content,
        mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        file_extension="xlsx",
    )

    assert result.page_count == 2
    assert result.author == "Sheet Author"


def test_extract_document_metadata_prefers_existing_page_count() -> None:
    content = _build_pptx()

    result = extract_document_metadata(
        content,
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        file_extension="pptx",
        existing_page_count=5,
    )

    assert result.page_count == 5


def test_extract_document_metadata_rejects_empty_bytes() -> None:
    with pytest.raises(MetadataExtractionError, match="empty"):
        extract_document_metadata(
            b"",
            mime_type="application/pdf",
            file_extension="pdf",
        )
