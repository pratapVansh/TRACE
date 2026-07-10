from io import BytesIO

import pytest
from pptx import Presentation

from app.services.document_processing_exceptions import PptxTextExtractionError
from app.services.pptx_text_extraction import extract_pptx_text


def _build_pptx(*slide_texts: str) -> bytes:
    presentation = Presentation()
    for text in slide_texts:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = text
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = f"{text} body"
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def test_extract_pptx_text_reads_every_slide() -> None:
    result = extract_pptx_text(_build_pptx("Slide One", "Slide Two"))

    assert result.slide_count == 2
    assert len(result.slides) == 2
    assert "Slide One" in result.slides[0].text
    assert "Slide Two" in result.slides[1].text
    assert "Slide One" in result.full_text
    assert "Slide Two" in result.full_text


def test_extract_pptx_text_handles_empty_presentation() -> None:
    presentation = Presentation()
    buffer = BytesIO()
    presentation.save(buffer)

    result = extract_pptx_text(buffer.getvalue())

    assert result.slide_count == 0
    assert result.full_text == ""


def test_extract_pptx_text_rejects_empty_bytes() -> None:
    with pytest.raises(PptxTextExtractionError, match="empty"):
        extract_pptx_text(b"")


def test_extract_pptx_text_rejects_invalid_bytes() -> None:
    with pytest.raises(PptxTextExtractionError, match="Failed to open PPTX"):
        extract_pptx_text(b"not-a-pptx")
