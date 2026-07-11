from datetime import datetime
from pathlib import Path
from zipfile import BadZipFile, ZipFile

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.core.config import settings
from app.core.logging import logger
from app.models.document import Document
from app.models.document_version import DocumentVersion
from app.processing.base import BaseProcessor

SLIDE_SEPARATOR = "\n\n" + "-" * 40 + "\n"
TABLE_SEPARATOR = "-" * 32


class PptxProcessor(BaseProcessor):
    name = "pptx_processor"
    supported_extensions = frozenset({"pptx"})
    supported_mime_types = frozenset({
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    })

    async def extract_text(
        self,
        document: Document,
        version: DocumentVersion,
    ) -> str:
        file_path = self._resolve_path(version)
        logger.info("Opening PPTX document_id=%s path=%s", document.id, file_path)

        try:
            prs = Presentation(file_path)
        except Exception as exc:
            logger.error("Failed to open PPTX document_id=%s: %s", document.id, exc)
            return ""

        parts: list[str] = []

        for slide_index, slide in enumerate(prs.slides):
            parts.append(SLIDE_SEPARATOR)
            parts.append(f"Slide {slide_index + 1}")
            parts.append(SLIDE_SEPARATOR)

            self._extract_slide_content(slide, parts)

            if slide.has_notes_slide:
                notes = slide.notes_slide
                notes_text = notes.notes_text_frame.text.strip()
                if notes_text:
                    parts.append("")
                    parts.append("Notes:")
                    parts.append(notes_text)

        return "\n".join(parts).strip()

    async def extract_metadata(
        self,
        document: Document,
        version: DocumentVersion,
    ) -> dict:
        file_path = self._resolve_path(version)

        try:
            prs = Presentation(file_path)
        except Exception as exc:
            logger.error("Failed to open PPTX for metadata document_id=%s: %s", document.id, exc)
            return self._fallback_metadata(document, version, file_path)

        cp = prs.core_properties
        slide_count = len(prs.slides)
        notes_count = 0
        table_count = 0
        total_text_boxes = 0
        image_count = 0
        chart_count = 0
        smartart_count = 0
        media_count = 0

        for slide in prs.slides:
            if slide.has_notes_slide:
                notes = slide.notes_slide
                if notes.notes_text_frame.text.strip():
                    notes_count += 1

            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table_count += 1
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_count += 1
                elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    chart_count += 1
                elif shape.shape_type == MSO_SHAPE_TYPE.DIAGRAM:
                    smartart_count += 1
                elif shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                    media_count += 1
                elif shape.shape_type in (MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.PLACEHOLDER):
                    total_text_boxes += 1
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    ic, cc, sc = self._count_group_shapes(shape)
                    image_count += ic
                    chart_count += cc
                    smartart_count += sc

        created = self._format_datetime(cp.created)
        modified = self._format_datetime(cp.modified)

        result = {
            "title": cp.title or document.title,
            "author": cp.author or "",
            "subject": cp.subject or "",
            "keywords": cp.keywords or "",
            "company": cp.company or "" if hasattr(cp, "company") else "",
            "created": created,
            "modified": modified,
            "last_modified_by": cp.last_modified_by or "",
            "revision": cp.revision,
            "slide_count": slide_count,
            "table_count": table_count,
            "image_count": image_count,
            "chart_count": chart_count,
            "notes_count": notes_count,
            "text_box_count": total_text_boxes,
            "smartart_count": smartart_count,
            "media_count": media_count,
            "file_size": self._get_file_size(file_path),
            "extension": "pptx",
            "mime_type": version.mime_type,
        }

        flags = []
        if image_count > 0:
            flags.append("image")
            result["requires_image_processing"] = True
        if chart_count > 0:
            flags.append("chart")
            result["requires_chart_processing"] = True
        if smartart_count > 0:
            flags.append("smartart")
            result["requires_smartart_processing"] = True
        if media_count > 0:
            flags.append("media")
            result["requires_media_processing"] = True

        if flags:
            logger.info(
                "Detected in document_id=%s: %s", document.id, ", ".join(flags)
            )

        logger.info(
            "Metadata extracted document_id=%s slides=%d tables=%d images=%d charts=%d notes=%d",
            document.id,
            slide_count,
            table_count,
            image_count,
            chart_count,
            notes_count,
        )
        return result

    async def validate(
        self,
        document: Document,
        version: DocumentVersion,
    ) -> list[str]:
        warnings: list[str] = []
        file_path = self._resolve_path(version)

        if self._is_encrypted_zip(file_path):
            warnings.append("Password-protected PPTX file")
            return warnings

        try:
            prs = Presentation(file_path)
        except BadZipFile:
            warnings.append("Corrupted or invalid PPTX file (not a valid ZIP archive)")
            return warnings
        except Exception as exc:
            err_str = str(exc).lower()
            if "password" in err_str or "encrypted" in err_str:
                warnings.append("Password-protected PPTX file")
            else:
                warnings.append(f"Cannot open PPTX: {exc}")
            return warnings

        if len(prs.slides) == 0:
            warnings.append("PPTX presentation is empty (no slides)")
            return warnings

        all_empty = True
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    all_empty = False
                    break
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                all_empty = False
                                break
                        if not all_empty:
                            break
                    if not all_empty:
                        break
                if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    all_empty = False
                    break
            if slide.has_notes_slide:
                notes = slide.notes_slide
                if notes.notes_text_frame.text.strip():
                    all_empty = False
            if not all_empty:
                break

        if all_empty:
            warnings.append("PPTX presentation contains no text content")

        return warnings

    async def process(
        self,
        document: Document,
        version: DocumentVersion,
    ) -> "ProcessingResult":
        from app.processing.models import ProcessingResult

        warnings = await self.validate(document, version)
        if warnings:
            metadata = await self.extract_metadata(document, version)
            return ProcessingResult(
                success=True,
                document_id=document.id,
                extracted_text="",
                metadata=metadata,
                warnings=warnings,
            )

        metadata = await self.extract_metadata(document, version)
        text = await self.extract_text(document, version)

        warnings_list: list[str] = []
        if metadata.get("requires_image_processing"):
            warnings_list.append("Embedded images detected. Image processing may be required.")
        if metadata.get("requires_chart_processing"):
            warnings_list.append("Charts detected. Chart data extraction not yet implemented.")
        if metadata.get("requires_smartart_processing"):
            warnings_list.append("SmartArt detected. SmartArt extraction not yet implemented.")
        if metadata.get("requires_media_processing"):
            warnings_list.append("Embedded media detected. Media extraction not yet implemented.")

        return ProcessingResult(
            success=True,
            document_id=document.id,
            extracted_text=text,
            metadata=metadata,
            warnings=warnings_list,
        )

    def _resolve_path(self, version: DocumentVersion) -> str:
        return str(settings.storage_root_path / version.storage_uri)

    def _extract_slide_content(self, slide, parts: list[str]) -> None:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                self._extract_table(shape.table, parts)

            elif shape.shape_type in (
                MSO_SHAPE_TYPE.TEXT_BOX,
                MSO_SHAPE_TYPE.PLACEHOLDER,
            ):
                if shape.has_text_frame:
                    self._extract_text_frame(shape.text_frame, parts)

            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self._extract_group_shapes(shape, parts)

            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if shape.has_text_frame:
                    caption = shape.text_frame.text.strip()
                    if caption:
                        parts.append(f"[Image: {caption}]")

            elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                if shape.has_chart:
                    chart = shape.chart
                    has_title = bool(chart.has_title and chart.chart_title and chart.chart_title.text_frame)
                    title_text = chart.chart_title.text_frame.text.strip() if has_title else "Untitled"
                    parts.append(f"[Chart: {title_text}]")

            elif shape.shape_type == MSO_SHAPE_TYPE.DIAGRAM:
                parts.append("[SmartArt Diagram]")

            elif shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                parts.append("[Embedded Media]")

    def _extract_text_frame(self, tf, parts: list[str]) -> None:
        lines: list[str] = []
        for paragraph in tf.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue

            level = paragraph.level if paragraph.level else 0
            indent = "  " * level

            is_bullet = self._is_bullet_paragraph(paragraph)
            is_numbered = not is_bullet and self._is_numbered_paragraph(paragraph)

            if is_bullet:
                lines.append(f"{indent}- {text}")
            elif is_numbered:
                lines.append(f"{indent}1. {text}")
            else:
                lines.append(f"{indent}{text}")

        if lines:
            parts.append("")
            parts.extend(lines)
            parts.append("")

    def _is_bullet_paragraph(self, paragraph) -> bool:
        pPr = paragraph._pPr
        if pPr is None:
            return False
        buChar = pPr.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}buChar")
        buNone = pPr.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}buNone")
        if buChar:
            return True
        if buNone:
            return False
        buFont = pPr.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}buFont")
        if buFont:
            return True
        return False

    def _is_numbered_paragraph(self, paragraph) -> bool:
        pPr = paragraph._pPr
        if pPr is None:
            return False
        buAutoNum = pPr.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum")
        return bool(buAutoNum)

    def _extract_table(self, table, parts: list[str]) -> None:
        parts.append("")
        parts.append(TABLE_SEPARATOR)

        rows = list(table.rows)
        if not rows:
            parts.append("(empty table)")
            parts.append(TABLE_SEPARATOR)
            parts.append("")
            return

        col_widths: list[int] = []
        for row in rows:
            for ci, cell in enumerate(row.cells):
                text = cell.text.strip()
                if ci >= len(col_widths):
                    col_widths.append(0)
                col_widths[ci] = max(col_widths[ci], len(text))

        col_widths = [max(w, 3) for w in col_widths]
        separator = " | ".join("-" * w for w in col_widths)

        for ri, row in enumerate(rows):
            cells = [cell.text.strip() for cell in row.cells]
            padded = []
            for ci, text in enumerate(cells):
                width = col_widths[ci] if ci < len(col_widths) else 20
                padded.append(text.ljust(width))
            parts.append(" | ".join(padded))

            if ri == 0:
                parts.append(separator)

        parts.append(TABLE_SEPARATOR)
        parts.append("")

    def _extract_group_shapes(self, group_shape, parts: list[str]) -> None:
        for shape in group_shape.shapes:
            if shape.has_text_frame:
                self._extract_text_frame(shape.text_frame, parts)
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                self._extract_table(shape.table, parts)
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if shape.has_text_frame:
                    caption = shape.text_frame.text.strip()
                    if caption:
                        parts.append(f"[Image: {caption}]")
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self._extract_group_shapes(shape, parts)

    def _count_group_shapes(self, group_shape) -> tuple[int, int, int]:
        ic: list[int] = [0]
        cc: list[int] = [0]
        sc: list[int] = [0]
        self._count_group_shapes_internal(group_shape, ic, cc, sc)
        return ic[0], cc[0], sc[0]

    def _count_group_shapes_internal(self, group_shape, ic, cc, sc) -> None:
        for shape in group_shape.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                ic[0] += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                cc[0] += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.DIAGRAM:
                sc[0] += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self._count_group_shapes_internal(shape, ic, cc, sc)

    @staticmethod
    def _is_encrypted_zip(file_path: str) -> bool:
        try:
            with ZipFile(file_path, "r") as zf:
                names = [n.lower() for n in zf.namelist()]
                for name in names:
                    if "encryptioninfo" in name:
                        return True
            return False
        except Exception:
            return False

    @staticmethod
    def _format_datetime(dt: datetime | None) -> str | None:
        if dt is None:
            return None
        try:
            return dt.isoformat()
        except (ValueError, AttributeError):
            return None

    @staticmethod
    def _get_file_size(file_path: str) -> int | None:
        try:
            return Path(file_path).stat().st_size
        except OSError:
            return None

    @staticmethod
    def _fallback_metadata(
        document: Document,
        version: DocumentVersion,
        file_path: str,
    ) -> dict:
        return {
            "title": document.title,
            "author": None,
            "slide_count": 0,
            "table_count": 0,
            "image_count": 0,
            "chart_count": 0,
            "notes_count": 0,
            "text_box_count": 0,
            "smartart_count": 0,
            "media_count": 0,
            "file_size": PptxProcessor._get_file_size(file_path),
            "extension": "pptx",
            "mime_type": version.mime_type,
        }
