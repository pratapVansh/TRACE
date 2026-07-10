from dataclasses import dataclass
from io import BytesIO

from pptx import Presentation

from app.services.document_processing_exceptions import PptxTextExtractionError

EXTRACTION_METHOD = "python-pptx"


@dataclass(frozen=True, slots=True)
class ExtractedSlide:
    slide_number: int
    text: str


@dataclass(frozen=True, slots=True)
class PptxTextExtractionResult:
    slides: tuple[ExtractedSlide, ...]
    full_text: str
    slide_count: int


def extract_pptx_text(content: bytes) -> PptxTextExtractionResult:
    """Extract text from every slide in a PPTX file."""
    if not content:
        raise PptxTextExtractionError("PPTX file is empty")

    try:
        presentation = Presentation(BytesIO(content))
    except Exception as exc:
        raise PptxTextExtractionError("Failed to open PPTX for text extraction") from exc

    slides: list[ExtractedSlide] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    parts.append(text)

        slides.append(
            ExtractedSlide(
                slide_number=slide_number,
                text="\n".join(parts),
            ),
        )

    full_text = "\n\n".join(slide.text for slide in slides if slide.text)
    return PptxTextExtractionResult(
        slides=tuple(slides),
        full_text=full_text,
        slide_count=len(slides),
    )
