#!/usr/bin/env python3
"""Generate the complete ABC Petrochemicals demo industrial dataset."""

import os
from datetime import datetime

OUTPUT = os.path.dirname(os.path.abspath(__file__))

# ── Consistent Data ──────────────────────────────────────────────
PERSONNEL = {
    "director":     ("Rajesh Mehta",      "Plant Director"),
    "it_manager":   ("Amit Patel",        "IT & Systems Manager"),
    "process_eng":  ("Priya Sharma",      "Senior Process Engineer"),
    "shift_lead":   ("Vikram Singh",      "Shift Lead Operator"),
    "compliance":   ("Ananya Gupta",      "Environmental Compliance Specialist"),
    "maint_eng":    ("Suresh Kumar",      "Maintenance Engineer"),
    "safety":       ("Deepak Joshi",      "Safety Officer"),
    "analyst":      ("Meena Iyer",        "Quality Analyst"),
}

EQUIPMENT = {
    "P-101": ("Centrifugal Pump",   "Cooling Water Circulation Pump"),
    "P-102": ("Centrifugal Pump",   "Condensate Return Pump"),
    "C-201": ("Air Compressor",     "Plant Instrument Air Compressor"),
    "B-101": ("Steam Boiler",       "High Pressure Steam Boiler"),
    "E-301": ("Heat Exchanger",     "Feed Preheater Shell & Tube"),
    "T-501": ("Storage Tank",       "Raw Material Storage Tank"),
    "CT-01": ("Cooling Tower",      "Induced Draft Cooling Tower"),
    "V-220": ("Control Valve",      "Steam Pressure Control Valve"),
    "M-110": ("Electric Motor",     "HVAC Drive Motor"),
    "L-401": ("Process Pipeline",   "Ethylene Transfer Pipeline"),
}

def make_output(subdir, filename):
    os.makedirs(os.path.join(OUTPUT, subdir), exist_ok=True)
    return os.path.join(OUTPUT, subdir, filename)

# ══════════════════════════════════════════════════════════════════
#  DOCX HELPERS
# ══════════════════════════════════════════════════════════════════

def _docx_header(doc, doc_id, title):
    from docx.shared import Pt
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "ABC Petrochemicals Pvt. Ltd."
    table.cell(0, 1).text = f"Document: {doc_id}"
    table.cell(1, 0).text = "Gujarat Manufacturing Plant"
    table.cell(1, 1).text = f"Rev: 1.0  |  Date: 15-Jan-2026"
    doc.add_heading(title, level=1)

def _docx_footer(doc, person_key):
    name, title = PERSONNEL[person_key]
    doc.add_paragraph("")
    p = doc.add_paragraph()
    p.add_run(f"Prepared By: {name}").bold = True
    doc.add_paragraph(f"Designation: {title}")
    doc.add_paragraph(f"Signature: ___________________")
    doc.add_paragraph("")
    p2 = doc.add_paragraph()
    p2.add_run("--- End of Document ---").italic = True

# ══════════════════════════════════════════════════════════════════
#  1. SOPs (3)
# ══════════════════════════════════════════════════════════════════

def generate_sops():
    from docx import Document
    from docx.shared import Pt
    files = []

    doc = Document()
    style = doc.styles["Normal"]; style.font.name = "Arial"; style.font.size = Pt(11)
    _docx_header(doc, "SOP-001", "Centrifugal Pump Start-Up Procedure")
    doc.add_paragraph("Equipment: P-101  --  Cooling Water Circulation Pump (Centrifugal Pump)")
    doc.add_paragraph("Department: Operations")
    doc.add_paragraph("Prepared By: Priya Sharma  |  Approved By: Rajesh Mehta\n")

    doc.add_heading("1. Purpose", level=2)
    doc.add_paragraph("To define the safe and systematic procedure for starting centrifugal pump P-101 to ensure reliable operation and prevent equipment damage.")
    doc.add_heading("2. Prerequisites", level=2)
    for s in ["Valid work permit obtained","Suction line fully open and vented","Electrical supply verified (415V, 3-phase)","Seal cooling water supply ON","Coupling guard securely fitted"]:
        doc.add_paragraph(s, style="List Bullet")
    doc.add_heading("3. Start-Up Procedure", level=2)
    steps = [
        "Verify lubrication oil level in bearing housing at sight-glass centerline.",
        "Open suction valve fully; confirm positive suction pressure.",
        "Open casing vent valve to release trapped air; then close.",
        "Start drive motor from DCS or local push-button station.",
        "Monitor discharge pressure rise to 4.5 barg within 5 seconds.",
        "Slowly open discharge valve to 25%, then gradually to 100%.",
        "Check motor current draw  --  must not exceed nameplate FLA (45 A).",
        "Inspect mechanical seal  --  minor weepage acceptable during settling.",
        "Verify pump vibration < 4.5 mm/s (alert limit).",
        "Log all start-up parameters in shift log.",
    ]
    for i, s in enumerate(steps, 1):
        doc.add_paragraph(f"{i}. {s}")
    doc.add_heading("4. Safety Precautions", level=2)
    for s in ["Ensure pump area is clear before starting.","Never start pump with suction valve closed.","Verify coupling guard is in place before energizing.","If abnormal vibration or noise occurs, stop immediately."]:
        doc.add_paragraph(s, style="List Bullet")
    doc.add_heading("5. Revision History", level=2)
    t = doc.add_table(rows=2, cols=4); t.style = "Light Shading Accent 1"
    for i, h in enumerate(["Rev","Date","Description","Author"]): t.cell(0,i).text = h
    for j, v in enumerate(["1.0","15-Jan-2026","Initial Release","Priya Sharma"]): t.cell(1,j).text = v
    doc.save(make_output("SOPs", "SOP-001_Pump_Start-Up_Procedure.docx"))
    files.append("SOPs/SOP-001_Pump_Start-Up_Procedure.docx")

    # SOP-002 Pump Shutdown
    doc = Document(); style = doc.styles["Normal"]; style.font.name = "Arial"; style.font.size = Pt(11)
    _docx_header(doc, "SOP-002", "Centrifugal Pump Shut-Down Procedure")
    doc.add_paragraph("Equipment: P-102  --  Condensate Return Pump (Centrifugal Pump)")
    doc.add_paragraph("Department: Operations")
    doc.add_paragraph("Prepared By: Suresh Kumar  |  Approved By: Rajesh Mehta\n")
    doc.add_heading("1. Purpose", level=2)
    doc.add_paragraph("To define the safe shut-down sequence for pump P-102 to prevent damage and ensure safe isolation for maintenance.")
    doc.add_heading("2. Procedure", level=2)
    steps = [
        "Notify production supervisor before shutdown.",
        "Confirm alternate pump is online if continuous operation required.",
        "Close discharge valve slowly to 10% to reduce backflow.",
        "Press STOP on local control station or DCS.",
        "Monitor pump coast-down  --  confirm smooth deceleration.",
        "Close discharge valve fully, then close suction valve.",
        "Open casing drain valve to relieve internal pressure.",
        "Close seal cooling water supply.",
        "Apply LOTO on the motor disconnect switch.",
        "Record shutdown time and reason in shift log.",
    ]
    for i, s in enumerate(steps, 1): doc.add_paragraph(f"{i}. {s}")
    doc.add_heading("3. Safety", level=2)
    for s in ["Hot surfaces  --  allow pump to cool before touching.","Verify LOTO before any maintenance.","Crack drain valve slowly; beware of trapped pressure."]:
        doc.add_paragraph(s, style="List Bullet")
    doc.save(make_output("SOPs", "SOP-002_Pump_Shut-Down_Procedure.docx"))
    files.append("SOPs/SOP-002_Pump_Shut-Down_Procedure.docx")

    # SOP-003 Boiler Operation
    doc = Document(); style = doc.styles["Normal"]; style.font.name = "Arial"; style.font.size = Pt(11)
    _docx_header(doc, "SOP-003", "Steam Boiler Start-Up Procedure")
    doc.add_paragraph("Equipment: B-101  --  High Pressure Steam Boiler")
    doc.add_paragraph("Department: Operations")
    doc.add_paragraph("Prepared By: Vikram Singh  |  Approved By: Rajesh Mehta\n")
    doc.add_heading("1. Purpose", level=2)
    doc.add_paragraph("To define the safe start-up sequence for HP steam boiler B-101 to achieve stable combustion and avoid thermal shock.")
    doc.add_heading("2. Pre-Start Checks", level=2)
    for s in ["Feed water quality: pH 8.5 - 9.5, conductivity < 10 µS/cm","Fuel gas supply pressure: 2.5 barg","Draft fan and FD fan tested","Gas leak check completed","All safety interlocks functional"]:
        doc.add_paragraph(s, style="List Bullet")
    doc.add_heading("3. Start-Up Sequence", level=2)
    steps = [
        "Verify drum water level at 50% (±25 mm) via sight glass.",
        "Open superheater vent valve for air purging.",
        "Start FD fan; purge furnace for 5 min at >= 30% air flow.",
        "Initiate pilot burner ignition; confirm flame scanner detection.",
        "Open main fuel gas valve; ramp main burner to 10% load.",
        "Increase firing rate max 5%/minute to avoid thermal stress.",
        "Close superheater vent at 1 barg drum pressure.",
        "Warm steam line by cracking main steam stop valve.",
        "Bring boiler to 42 barg at <= 50 deg C/hr heating rate.",
        "Place boiler in AUTO mode once stable.",
    ]
    for i, s in enumerate(steps, 1): doc.add_paragraph(f"{i}. {s}")
    doc.add_heading("4. Safety", level=2)
    for s in ["Never start without completing furnace purge.","Monitor drum water level constantly.","Adhere to ramp rates to prevent thermal shock.","Two operators must be present during start-up."]:
        doc.add_paragraph(s, style="List Bullet")
    doc.save(make_output("SOPs", "SOP-003_Boiler_Start-Up_Procedure.docx"))
    files.append("SOPs/SOP-003_Boiler_Start-Up_Procedure.docx")
    return files

# ══════════════════════════════════════════════════════════════════
#  2. OEM Manuals (2)
# ══════════════════════════════════════════════════════════════════

def generate_manuals():
    from docx import Document; from docx.shared import Pt
    files = []

    doc = Document(); style = doc.styles["Normal"]; style.font.name = "Arial"; style.font.size = Pt(11)
    _docx_header(doc, "MAN-001", "Centrifugal Pump OEM Manual  --  P-101 / P-102")
    doc.add_paragraph("Equipment: P-101, P-102  --  Centrifugal Pumps")
    doc.add_paragraph("Supplier: Kirloskar Brothers Ltd.")
    doc.add_paragraph("Compiled By: Priya Sharma  |  Date: 20-Dec-2025\n")
    doc.add_heading("1. General Description", level=2)
    doc.add_paragraph("The KBL Series-C horizontal centrifugal pump is designed for continuous industrial service handling clean or slightly contaminated liquids. Models P-101 and P-102 are identical frame sizes with different impeller trims.")
    doc.add_heading("2. Technical Specifications", level=2)
    t = doc.add_table(rows=8, cols=3); t.style = "Light Shading Accent 1"
    for i, h in enumerate(["Parameter","P-101","P-102"]):
        for j, v in enumerate([h,"Flow Rate","Head","Speed","Motor Power","Impeller Dia.","Material","Connection"]):
            if i == 0: t.cell(j,i).text = v
    data = ["Parameter","P-101","P-102","Flow Rate","180 m^3/hr","120 m^3/hr","Head","65 m","50 m","Speed","2950 RPM","2950 RPM","Motor Power","45 kW","30 kW","Impeller Dia.","240 mm","210 mm","Material","CI / SS316","CI / SS304","Connection","6\" Flanged","4\" Flanged"]
    for idx in range(0, len(data), 3):
        r = idx//3
        for c in range(3):
            t.cell(r,c).text = data[idx+c]
    doc.add_heading("3. Installation", level=2)
    doc.add_paragraph("Foundation must be level within 0.05 mm/m. Grout after alignment. Suction piping must be straight for minimum 5 pipe diameters before pump suction. Use pipe supports to avoid nozzle loading.")
    doc.add_heading("4. Maintenance Schedule", level=2)
    doc.add_paragraph("• Daily: Check bearing temperature, seal leakage, vibration")
    doc.add_paragraph("• Monthly: Grease bearings, check coupling alignment")
    doc.add_paragraph("• Quarterly: Check impeller clearance, replace if worn")
    doc.add_paragraph("• Annually: Major overhaul  --  replace bearings, seals, gaskets")
    doc.add_heading("5. Troubleshooting", level=2)
    doc.add_paragraph("Low Flow: Check suction strainer, impeller clearance, NPSH available.")
    doc.add_paragraph("High Vibration: Check alignment, bearing condition, cavitation.")
    doc.add_paragraph("Seal Leakage: Replace mechanical seal; verify seal water quality.")
    doc.save(make_output("Manuals", "MAN-001_Centrifugal_Pump_Manual.docx"))
    files.append("Manuals/MAN-001_Centrifugal_Pump_Manual.docx")

    doc = Document(); style = doc.styles["Normal"]; style.font.name = "Arial"; style.font.size = Pt(11)
    _docx_header(doc, "MAN-002", "Air Compressor OEM Manual  --  C-201")
    doc.add_paragraph("Equipment: C-201  --  Plant Instrument Air Compressor")
    doc.add_paragraph("Supplier: Atlas Copco (Model: ZR-160)")
    doc.add_paragraph("Compiled By: Suresh Kumar  |  Date: 18-Dec-2025\n")
    doc.add_heading("1. General Description", level=2)
    doc.add_paragraph("The Atlas Copco ZR-160 is an oil-free rotary screw compressor designed for instrument air supply. It delivers 160 m^3/hr at 8.5 barg with integrated dryer and filtration.")
    doc.add_heading("2. Specifications", level=2)
    for kv in [("Type","Oil-free rotary screw"),("Flow","160 m^3/hr FAD"),("Pressure","8.5 barg"),("Motor","200 kW, 415V, 3-phase"),("Cooling","Water-cooled"),("Dryer","Heatless regenerative"),("Dew Point","-40 deg C PDP")]:
        doc.add_paragraph(f"{kv[0]}: {kv[1]}")
    doc.add_heading("3. Operating Parameters", level=2)
    doc.add_paragraph("• Discharge temperature: 85 - 110 deg C (alarm at 120 deg C, trip at 130 deg C)")
    doc.add_paragraph("• Oil pressure: 3.0 - 4.5 barg")
    doc.add_paragraph("• Cooling water inlet: max 32 deg C")
    doc.add_paragraph("• Vibration: < 7.1 mm/s (alert at 11 mm/s)")
    doc.add_heading("4. Lubrication", level=2)
    doc.add_paragraph("Use Atlas Copco Oil Roto-Xtend ND. Oil change interval: 4000 hours. Grease motor bearings every 2000 hours with SKF LGHP 2.")
    doc.add_heading("5. Safety Devices", level=2)
    doc.add_paragraph("• Pressure safety valve set at 10.5 barg")
    doc.add_paragraph("• Motor thermal overload relay")
    doc.add_paragraph("• High discharge temperature trip")
    doc.add_paragraph("• Emergency stop push-button at local panel")
    doc.save(make_output("Manuals", "MAN-002_Air_Compressor_Manual.docx"))
    files.append("Manuals/MAN-002_Air_Compressor_Manual.docx")
    return files

# ══════════════════════════════════════════════════════════════════
#  3. Maintenance Reports (3)
# ══════════════════════════════════════════════════════════════════

def generate_maintenance():
    from docx import Document; from docx.shared import Pt
    files = []

    doc = Document(); style = doc.styles["Normal"]; style.font.name = "Arial"; style.font.size = Pt(11)
    _docx_header(doc, "MNT-001", "Quarterly Preventive Maintenance  --  Cooling Tower CT-01")
    doc.add_paragraph("Equipment: CT-01  --  Induced Draft Cooling Tower")
    doc.add_paragraph("Date: 05-Apr-2026")
    doc.add_paragraph("Performed By: Suresh Kumar (Maintenance Engineer)\n")
    doc.add_heading("Work Summary", level=2)
    doc.add_paragraph("Routine quarterly preventive maintenance completed on cooling tower CT-01. All six cells inspected.")
    doc.add_heading("Activities Performed", level=2)
    for s in ["Inspected fan gearbox oil level and condition  --  OK","Greased fan shaft bearings (SKF LGHP 2)","Cleaned and flushed water distribution nozzles","Checked and tensioned V-belts on all 4 fans per cell","Inspected fill media for fouling and scaling","Measured vibration on fan gearbox  --  max 3.2 mm/s","Water treatment chemical dosing verified","Calibrated conductivity controller"]:
        doc.add_paragraph(s, style="List Bullet")
    doc.add_heading("Findings", level=2)
    doc.add_paragraph("Cell #3 fill media shows moderate scaling on the hot side. Recommended chemical cleaning during next scheduled turnaround.")
    doc.add_heading("Recommendations", level=2)
    doc.add_paragraph("1. Schedule fill media chemical cleaning for Q3 2026.")
    doc.add_paragraph("2. Increase biocide dosing frequency during monsoon months.")
    doc.add_heading("Sign-Off", level=2)
    doc.add_paragraph("Maintained By: Suresh Kumar  |  Signature: sk/ct01/05apr26")
    doc.save(make_output("Maintenance", "MNT-001_Quarterly_PM_Cooling_Tower.docx"))
    files.append("Maintenance/MNT-001_Quarterly_PM_Cooling_Tower.docx")

    doc = Document(); style = doc.styles["Normal"]; style.font.name = "Arial"; style.font.size = Pt(11)
    _docx_header(doc, "MNT-002", "Bearing Replacement Report  --  P-101")
    doc.add_paragraph("Equipment: P-101  --  Cooling Water Circulation Pump")
    doc.add_paragraph("Date: 12-Mar-2026")
    doc.add_paragraph("Work Order: WO-2026-0341\n")
    doc.add_heading("Reason for Replacement", level=2)
    doc.add_paragraph("High vibration (7.8 mm/s) and elevated bearing temperature (78 deg C) observed during routine monitoring. Investigation indicated bearing wear on the drive-end (DE) bearing.")
    doc.add_heading("Work Performed", level=2)
    for s in ["Isolated P-101 and applied LOTO","Removed coupling guard and coupling halves","Extracted DE bearing (SKF 6319 C3) and NDE bearing (SKF 6320 C3)","Inspected bearing housing and shaft  --  no scoring","Installed new DE and NDE bearings using induction heater (max 110 deg C)","Reassembled coupling and aligned within 0.05 mm","Filled bearing housing with fresh grease (Mobil Polyrex EM)","Reinstalled coupling guard and removed LOTO","Performed test run for 2 hours"]:
        doc.add_paragraph(s, style="List Bullet")
    doc.add_heading("Post-Repair Readings", level=2)
    doc.add_paragraph("Vibration (DE): 1.8 mm/s")
    doc.add_paragraph("Vibration (NDE): 2.1 mm/s")
    doc.add_paragraph("Bearing Temperature: 52 deg C (stabilized)")
    doc.add_heading("Parts Used", level=2)
    doc.add_paragraph("SKF 6319 C3 (Qty 1)  --  Bearing No. B-6319-01")
    doc.add_paragraph("SKF 6320 C3 (Qty 1)  --  Bearing No. B-6320-01")
    doc.add_paragraph("Mobil Polyrex EM  --  Grease (500 g)")
    doc.add_heading("Sign-Off", level=2)
    doc.add_paragraph("Performed By: Suresh Kumar  |  Verified By: Priya Sharma")
    doc.save(make_output("Maintenance", "MNT-002_Bearing_Replacement_P-101.docx"))
    files.append("Maintenance/MNT-002_Bearing_Replacement_P-101.docx")

    doc = Document(); style = doc.styles["Normal"]; style.font.name = "Arial"; style.font.size = Pt(11)
    _docx_header(doc, "MNT-003", "Mechanical Seal Leakage Repair  --  P-102")
    doc.add_paragraph("Equipment: P-102  --  Condensate Return Pump")
    doc.add_paragraph("Date: 22-Feb-2026")
    doc.add_paragraph("Work Order: WO-2026-0287\n")
    doc.add_heading("Problem Description", level=2)
    doc.add_paragraph("Operator reported visible drip leakage (~15 drops/min) from the mechanical seal area of pump P-102. Seal water flush line was at normal pressure.")
    doc.add_heading("Root Cause", level=2)
    doc.add_paragraph("Mechanical seal faces worn beyond tolerance. Seal face exhibited circumferential scoring and one visible chip on the stationary face.")
    doc.add_heading("Repair Actions", level=2)
    for s in ["Isolated pump and performed LOTO","Drained pump casing and removed backplate","Extracted defective mechanical seal (John Crane 5610, 45mm)","Inspected shaft sleeve  --  minor wear marks, polished with fine emery","Installed new mechanical seal (John Crane 5610, 45mm) with clean assembly gel","Reassembled backplate and filled casing","Slowly pressurized and checked for leaks  --  zero leakage at 4.5 barg","Started pump and monitored for 1 hour  --  seal weepage < 3 drops/min"]:
        doc.add_paragraph(s, style="List Bullet")
    doc.add_heading("Parts Used", level=2)
    doc.add_paragraph("John Crane 5610 Mechanical Seal, 45mm (Qty 1)  --  Part No. JC-5610-45")
    doc.add_paragraph("Gasket set (pump backplate)  --  P-102-GSKT-01")
    doc.add_heading("Recommendation", level=2)
    doc.add_paragraph("Monitor seal water supply quality. Install Y-strainer on seal water line to prevent particulate damage to seal faces.")
    doc.add_heading("Sign-Off", level=2)
    doc.add_paragraph("Performed By: Suresh Kumar  |  Approved By: Priya Sharma")
    doc.save(make_output("Maintenance", "MNT-003_Seal_Leakage_Repair_P-102.docx"))
    files.append("Maintenance/MNT-003_Seal_Leakage_Repair_P-102.docx")
    return files

# ══════════════════════════════════════════════════════════════════
#  4. Inspection Reports (3)
# ══════════════════════════════════════════════════════════════════

def generate_inspections():
    from docx import Document; from docx.shared import Pt
    files = []

    doc = Document(); style = doc.styles["Normal"]; style.font.name = "Arial"; style.font.size = Pt(11)
    _docx_header(doc, "INS-001", "Pressure Vessel Inspection  --  T-501")
    doc.add_paragraph("Equipment: T-501  --  Raw Material Storage Tank")
    doc.add_paragraph("Inspection Date: 18-Jun-2026")
    doc.add_paragraph("Inspector: Deepak Joshi (Safety Officer)\n")
    doc.add_heading("Inspection Type", level=2)
    doc.add_paragraph("External visual inspection per API 653  --  5-year interval.")
    doc.add_heading("Observations", level=2)
    for s in ["External coating intact  --  no significant rust or peeling noted at accessible areas","Nameplate legible; last inspection date confirmed as June 2021","Stairway and platform in good condition; handrails compliant","Foundation  --  no visible settlement or cracking","Grounding cable resistance measured: 2.3 Ω (limit: < 10 Ω)","Vent piping clear; no obstruction","Overflow pipe unrestricted"]:
        doc.add_paragraph(s, style="List Bullet")
    doc.add_heading("UT Spot Check Results", level=2)
    doc.add_paragraph("Bottom plate (4 spots): 7.2, 7.5, 7.0, 7.3 mm (min req'd: 6.0 mm)")
    doc.add_paragraph("Shell course 1 (4 spots): 9.8, 10.1, 9.9, 10.0 mm (min req'd: 8.0 mm)")
    doc.add_heading("Conclusion", level=2)
    doc.add_paragraph("Tank T-501 is in satisfactory condition. All measured thicknesses exceed minimum required values. Next full inspection due: June 2031.")
    doc.add_paragraph("Inspector: Deepak Joshi  |  Signature: dj/t501/18jun26")
    doc.save(make_output("Inspection", "INS-001_Pressure_Vessel_Inspection_T-501.docx"))
    files.append("Inspection/INS-001_Pressure_Vessel_Inspection_T-501.docx")

    doc = Document(); style = doc.styles["Normal"]; style.font.name = "Arial"; style.font.size = Pt(11)
    _docx_header(doc, "INS-002", "Vibration Analysis Report  --  C-201")
    doc.add_paragraph("Equipment: C-201  --  Plant Instrument Air Compressor")
    doc.add_paragraph("Survey Date: 10-May-2026")
    doc.add_paragraph("Analyst: Suresh Kumar (Maintenance Engineer)\n")
    doc.add_heading("Methodology", level=2)
    doc.add_paragraph("Vibration data collected using SKF Microlog CMXA-75 analyzer. Measurements taken at 8 bearing points in X, Y, Z axes. Data compared against ISO 10816-3 criteria.")
    doc.add_heading("Measurement Summary", level=2)
    t = doc.add_table(rows=5, cols=4); t.style = "Light Shading Accent 1"
    for j, h in enumerate(["Point","X (mm/s)","Y (mm/s)","Z (mm/s)"]): t.cell(0,j).text = h
    pts = [("Motor DE",1.8,2.1,1.5),("Motor NDE",2.0,2.3,1.7),("Comp DE",3.5,4.1,2.9),("Comp NDE",3.2,3.8,2.6)]
    for r_idx, (pt,x,y,z) in enumerate(pts,1):
        t.cell(r_idx,0).text = pt; t.cell(r_idx,1).text=str(x); t.cell(r_idx,2).text=str(y); t.cell(r_idx,3).text=str(z)
    doc.add_heading("Analysis", level=2)
    doc.add_paragraph("Overall vibration levels at compressor bearings are in the ALERT zone per ISO 10816-3 (2.8 - 4.5 mm/s, Class II). The dominant frequency peak at 1x RPM (2950 CPM) suggests slight unbalance. No bearing defect frequencies detected.")
    doc.add_heading("Recommendations", level=2)
    doc.add_paragraph("1. Schedule dynamic balancing of compressor rotor during next planned shutdown.")
    doc.add_paragraph("2. Increase monitoring frequency to monthly until resolved.")
    doc.add_paragraph("3. Check foundation bolt torque during next weekly inspection.")
    doc.save(make_output("Inspection", "INS-002_Vibration_Analysis_C-201.docx"))
    files.append("Inspection/INS-002_Vibration_Analysis_C-201.docx")

    doc = Document(); style = doc.styles["Normal"]; style.font.name = "Arial"; style.font.size = Pt(11)
    _docx_header(doc, "INS-003", "Boiler Inspection Report  --  B-101")
    doc.add_paragraph("Equipment: B-101  --  High Pressure Steam Boiler")
    doc.add_paragraph("Inspection Date: 15-May-2026")
    doc.add_paragraph("Inspector: Deepak Joshi (Safety Officer)\n")
    doc.add_heading("Inspection Scope", level=2)
    doc.add_paragraph("Annual external inspection per IBR (Indian Boiler Regulations). Internal inspection deferred to next shutdown.")
    doc.add_heading("External Inspection Findings", level=2)
    for s in ["Furnace casing  --  no hot spots or bulging detected","Refractory condition  --  minor cracks in burner tile area (noted for repair)","Safety valves (2x)  --  lifting pressure tested at 44 barg and 44.5 barg; reset correctly","Water level gauges  --  clear and operational; blowdown valves functional","Feed water check valve  --  OK, non-return function verified","Insulation  --  intact on drum and headers"]:
        doc.add_paragraph(s, style="List Bullet")
    doc.add_heading("Non-Destructive Testing", level=2)
    doc.add_paragraph("UT thickness readings on steam drum (shell):")
    doc.add_paragraph("  Bottom: 32.5 mm / 31.8 mm / 32.1 mm (min req'd: 28.0 mm)")
    doc.add_paragraph("  Top:   34.2 mm / 33.9 mm / 34.0 mm (min req'd: 28.0 mm)")
    doc.add_heading("Actions Required", level=2)
    doc.add_paragraph("1. Repair burner tile refractory during next shutdown.")
    doc.add_paragraph("2. Internal inspection to be completed within 6 months.")
    doc.add_paragraph("Signed: Deepak Joshi  |  Witnessed: Vikram Singh")
    doc.save(make_output("Inspection", "INS-003_Boiler_Inspection_B-101.docx"))
    files.append("Inspection/INS-003_Boiler_Inspection_B-101.docx")
    return files

# ══════════════════════════════════════════════════════════════════
#  5. Incident Reports (2)
# ══════════════════════════════════════════════════════════════════

def generate_incidents():
    from docx import Document; from docx.shared import Pt
    files = []

    doc = Document(); style = doc.styles["Normal"]; style.font.name = "Arial"; style.font.size = Pt(11)
    _docx_header(doc, "INC-001", "Pump Oil Leakage Incident  --  P-101")
    doc.add_paragraph("Equipment: P-101  --  Cooling Water Circulation Pump")
    doc.add_paragraph("Incident Date: 03-Apr-2026 | Time: 14:30 hrs")
    doc.add_paragraph("Reported By: Vikram Singh (Shift Lead)\n")
    doc.add_heading("Incident Description", level=2)
    doc.add_paragraph('During routine round at 14:30 hrs, operator observed a steady oil leak from the pump P-101 bearing housing. Approximately 2 litres of lubricating oil had pooled on the drip tray. The pump was operating normally with no abnormal vibration or noise.')
    doc.add_heading("Immediate Actions Taken", level=2)
    for s in ["Pump was taken offline and spare pump (P-101-S) was started","Oil spill contained using absorbent pads  --  area cleaned","Bearing housing drain plug found to be loose  --  tightened to 35 Nm","Fresh oil (Mobil Polyrex EM) replenished to correct level","Pump restarted and monitored for 1 hour  --  no further leakage"]:
        doc.add_paragraph(s, style="List Bullet")
    doc.add_heading("Root Cause", level=2)
    doc.add_paragraph("The bearing housing drain plug was not fully tightened during the previous bearing replacement (MNT-002, 12-Mar-2026). Vibration during operation caused gradual loosening over 22 days of service.")
    doc.add_heading("Corrective Actions", level=2)
    doc.add_paragraph("1. Added drain plug tightening to bearing replacement checklist.")
    doc.add_paragraph("2. Lock wire applied to drain plug to prevent back-off.")
    doc.add_paragraph("3. Maintenance engineer (Suresh Kumar) counselled on torque verification.")
    doc.add_heading("Environmental Impact", level=2)
    doc.add_paragraph("Approximately 2 litres of oil contained within drip tray. No soil or water contamination. Absorbent pads disposed as hazardous waste.")
    doc.add_paragraph("Reported By: Vikram Singh  |  Reviewed By: Rajesh Mehta")
    doc.save(make_output("Incidents", "INC-001_Pump_Oil_Leakage_P-101.docx"))
    files.append("Incidents/INC-001_Pump_Oil_Leakage_P-101.docx")

    doc = Document(); style = doc.styles["Normal"]; style.font.name = "Arial"; style.font.size = Pt(11)
    _docx_header(doc, "INC-002", "Control Valve Failure  --  V-220")
    doc.add_paragraph("Equipment: V-220  --  Steam Pressure Control Valve")
    doc.add_paragraph("Incident Date: 18-Feb-2026 | Time: 09:15 hrs")
    doc.add_paragraph("Reported By: Vikram Singh (Shift Lead)\n")
    doc.add_heading("Incident Description", level=2)
    doc.add_paragraph('At 09:15 hrs, the DCS alarmed "Steam Header Pressure High." Investigation revealed control valve V-220 failed to respond to the control signal. The valve was stuck at 45% open while the controller demanded 25%. Downstream pressure rose to 8.2 barg (alarm at 7.5 barg). The emergency steam vent valve (PV-2201) automatically opened at 8.0 barg, preventing a full overpressure event.')
    doc.add_heading("Immediate Actions", level=2)
    for s in ["Placed V-220 in MANUAL mode  --  position did not respond","Isolated steam supply to valve (bypass valve opened manually)","Actuator air supply checked  --  4.2 barg (within spec)","Scheduled emergency maintenance (WO-2026-0251)"]:
        doc.add_paragraph(s, style="List Bullet")
    doc.add_heading("Root Cause", level=2)
    doc.add_paragraph("Valve actuator diaphragm ruptured due to age-related degradation. The actuator was 8 years in service. The diaphragm material (EPDM) showed signs of thermal cracking.")
    doc.add_heading("Corrective Actions", level=2)
    doc.add_paragraph("1. Replaced actuator diaphragm (Fisher part no. DIA-220-EPDM-01).")
    doc.add_paragraph("2. Bench-set calibration verified  --  valve strokes 0 - 100% within spec.")
    doc.add_paragraph("3. Added V-220 to critical valve PM schedule with 2-year diaphragm replacement interval.")
    doc.add_paragraph("Reported By: Vikram Singh  |  Reviewed By: Priya Sharma")
    doc.save(make_output("Incidents", "INC-002_Valve_Failure_V-220.docx"))
    files.append("Incidents/INC-002_Valve_Failure_V-220.docx")
    return files

# ══════════════════════════════════════════════════════════════════
#  6. Shift Logs (2 TXT)
# ══════════════════════════════════════════════════════════════════

def generate_shiftlogs():
    files = []
    log1 = """======================================================================
ABC Petrochemicals Pvt. Ltd.  --  Gujarat Manufacturing Plant
SHIFT LOG  --  MORNING SHIFT (06:00  -  14:00)
======================================================================
Date:      15-Jul-2026
Shift Lead: Vikram Singh
Operators:  Ramesh Kumar, Sunil Patil, Jitendra Verma

======================================================================
PRODUCTION SUMMARY
======================================================================
Ethylene Cracker   : 95% throughput (target: 98%)
Propylene Unit     : 100% throughput
Utilities          : All within limits

================================================================------
EQUIPMENT STATUS
======================================================================
P-101  : Running normally. Bearing temp 54 deg C, vibration 2.1 mm/s.
P-102  : Standby (isolated). Seal repair completed 22-Feb-2026.
C-201  : Running. Vibration monitored  --  see INS-002 for details.
B-101  : Running at 42 barg / 440 deg C. Boiler inspection due.
E-301  : Duty service. Outlet temp 118 deg C (normal).
T-501  : Level 65%. Receiving feed at 12 m^3/hr.
CT-01  : All 6 cells in service. Basin temp 31 deg C.
V-220  : 38% open. Manual mode until actuator replaced. [ACTIVE WORK]
M-110  : Running. Motor current 38 A (FLA: 45 A).
L-401  : Ethylene transfer at 85% rate. Pipeline pigging scheduled.

======================================================================
SHIFT EVENTS
======================================================================
08:45   --  Fire water pump run test completed (15 min run). OK.
09:30   --  Contractor toolbox talk: confined space entry.
11:00   --  E-301 feed pump seal water strainer cleaned.
12:15   --  Spare P-102 seal flush line flushed and checked.
13:00   --  Cooling tower water analysis: pH 7.8, Cl 85 ppm.

======================================================================
SAFETY OBSERVATIONS
======================================================================
  • One missing lockout tag on utility panel (replaced immediately).
  • Good housekeeping observed in cracker unit.

======================================================================
HANDOVER NOTES
======================================================================
  --> Night shift to monitor V-220 actuator replacement progress.
  --> T-501 inspection report (INS-001) to be filed.
  --> Keep P-101 standby pump (P-101-S) available.

Handover By: Vikram Singh    |    Time: 13:45
Received By: N/A (night shift)
======================================================================
"""
    fp = make_output("ShiftLogs", "LOG-001_Morning_Shift_15-Jul-2026.txt")
    with open(fp, "w", encoding="utf-8") as f: f.write(log1)
    files.append("ShiftLogs/LOG-001_Morning_Shift_15-Jul-2026.txt")

    log2 = """======================================================================
ABC Petrochemicals Pvt. Ltd.  --  Gujarat Manufacturing Plant
SHIFT LOG  --  NIGHT SHIFT (22:00  -  06:00)
======================================================================
Date:      15-Jul-2026 / 16-Jul-2026
Shift Lead: Sanjay Patel
Operators:  Ashok Yadav, Manoj Desai

======================================================================
EQUIPMENT STATUS
======================================================================
P-101  : Running normally. Vibration steady at 2.3 mm/s.
P-102  : Standby.
C-201  : Running. No alarms.
B-101  : Auto mode. Steam production steady at 140 t/hr.
E-301  : Normal duty. Temperatures stable.
T-501  : Level 62%. Night fill valve closed.
CT-01  : Normal. Basin level maintained.
V-220  : Replaced actuator diaphragm. Bench test OK. Valve returned to
         AUTO mode at 23:30 hrs. Stroke test passed  --  0% to 100% in 12 sec.
M-110  : Running. Temperature 62 deg C.
L-401  : Ethylene transfer rate steady at 85%. No pigging issues.

======================================================================
SHIFT EVENTS
======================================================================
22:30   --  V-220 actuator diaphragm replacement completed.
23:00   --  V-220 bench calibration and stroking test. Verified OK.
23:30   --  V-220 returned to service in AUTO mode. Downstream pressure: 7.2 barg.
01:15   --  Routine round: all equipment operating normally.
03:00   --  B-101 soot blowing cycle completed (manual mode).
04:30   --  T-501 level check and data logged.
05:00   --  Pre-day-shift preparation: log sheets updated.

======================================================================
SAFETY OBSERVATIONS
======================================================================
  • V-220 work area barricaded properly during maintenance.
  • All LOTO tags accounted for and removed after V-220 job completion.
  • No safety incidents during the shift.

======================================================================
HANDOVER NOTES
======================================================================
  --> V-220 actuator replacement complete and tested. Ready for normal service.
  --> Morning to monitor E-301 outlet temperature trend.
  --> CT-01 chemical cleaning to be planned for Q3.
  --> B-101 internal inspection due  --  coordinate with maintenance.

Handover By: Sanjay Patel    |    Time: 05:45
Received By: Vikram Singh (morning shift)
======================================================================
"""
    fp = make_output("ShiftLogs", "LOG-002_Night_Shift_15-16-Jul-2026.txt")
    with open(fp, "w", encoding="utf-8") as f: f.write(log2)
    files.append("ShiftLogs/LOG-002_Night_Shift_15-16-Jul-2026.txt")
    return files

# ══════════════════════════════════════════════════════════════════
#  7. Excel Files (3 XLSX)
# ══════════════════════════════════════════════════════════════════

def generate_excel():
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    files = []

    hdr_font = Font(bold=True, color="FFFFFF", size=11)
    hdr_fill = PatternFill(start_color="2F5496", end_color="2F5496", fill_type="solid")
    thin_border = Border(
        left=Side(style='thin'), right=Side(style='thin'),
        top=Side(style='thin'), bottom=Side(style='thin')
    )

    def style_header(ws, row=1, cols=None):
        for c in range(1, (cols or ws.max_column)+1):
            cell = ws.cell(row=row, column=c)
            cell.font = hdr_font; cell.fill = hdr_fill; cell.alignment = Alignment(horizontal="center")
            cell.border = thin_border

    # --- Equipment Register ---
    wb = Workbook(); ws = wb.active; ws.title = "Equipment Register"
    headers = ["Asset ID","Asset Name","Type","Department","Location","Install Date","Status","Criticality"]
    ws.append(headers); style_header(ws)
    data = [
        ["P-101","Cooling Water Circulation Pump","Centrifugal Pump","Operations","Cracker Unit","15-Mar-2015","Active","High"],
        ["P-102","Condensate Return Pump","Centrifugal Pump","Operations","Utilities","22-Jun-2018","Active","Medium"],
        ["C-201","Plant Instrument Air Compressor","Air Compressor","Maintenance","Compressor House","10-Jan-2020","Active","High"],
        ["B-101","High Pressure Steam Boiler","Steam Boiler","Operations","Boiler House","05-Nov-1998","Active","Critical"],
        ["E-301","Feed Preheater","Heat Exchanger","Engineering","Cracker Unit","30-Aug-2015","Active","High"],
        ["T-501","Raw Material Storage Tank","Storage Tank","Operations","Tank Farm","12-Mar-1999","Active","Critical"],
        ["CT-01","Induced Draft Cooling Tower","Cooling Tower","Utilities","Cooling Tower Yard","20-Jul-2010","Active","Medium"],
        ["V-220","Steam Pressure Control Valve","Control Valve","Technical","Steam Header","15-Jun-2018","Active","High"],
        ["M-110","HVAC Drive Motor","Electric Motor","Maintenance","HVAC Room","05-Sep-2019","Active","Low"],
        ["L-401","Ethylene Transfer Pipeline","Process Pipeline","Engineering","Pipe Rack","22-Jan-2005","Active","Critical"],
    ]
    for row_data in data:
        ws.append(row_data)
    for row in ws.iter_rows(min_row=2, max_row=ws.max_row, max_col=len(headers)):
        for cell in row:
            cell.border = thin_border
    ws.column_dimensions['A'].width = 12; ws.column_dimensions['B'].width = 38
    ws.column_dimensions['C'].width = 20; ws.column_dimensions['D'].width = 15
    ws.column_dimensions['E'].width = 20; ws.column_dimensions['F'].width = 14
    ws.column_dimensions['G'].width = 10; ws.column_dimensions['H'].width = 12
    wb.save(make_output("Excel", "Equipment_Register.xlsx"))
    files.append("Excel/Equipment_Register.xlsx")

    # --- Maintenance Schedule ---
    wb = Workbook(); ws = wb.active; ws.title = "Schedule"
    headers = ["Asset ID","Task Description","Frequency","Last Done","Next Due","Assigned To"]
    ws.append(headers); style_header(ws)
    data = [
        ["P-101","Bearing greasing","Monthly","01-Jul-2026","01-Aug-2026","Suresh Kumar"],
        ["P-101","Coupling alignment check","Quarterly","15-May-2026","15-Aug-2026","Suresh Kumar"],
        ["P-101","Major overhaul","Annual","12-Mar-2026","12-Mar-2027","Suresh Kumar"],
        ["P-102","Mechanical seal check","Quarterly","22-Feb-2026","22-May-2026","Suresh Kumar"],
        ["C-201","Oil change","4000 hrs","20-Jan-2026","20-Aug-2026","Suresh Kumar"],
        ["C-201","Dryer desiccant replacement","2 years","10-Jan-2024","10-Jan-2026","Suresh Kumar"],
        ["B-101","Safety valve testing","Annual","15-May-2026","15-May-2027","Deepak Joshi"],
        ["B-101","Internal inspection","5 years","15-May-2021","15-May-2026","Deepak Joshi"],
        ["E-301","Tube bundle cleaning","2 years","30-Aug-2024","30-Aug-2026","Suresh Kumar"],
        ["T-501","External visual inspection","5 years","18-Jun-2026","18-Jun-2031","Deepak Joshi"],
        ["CT-01","Fan gearbox oil change","Annual","05-Apr-2026","05-Apr-2027","Suresh Kumar"],
        ["V-220","Diaphragm replacement","2 years","15-Jul-2026","15-Jul-2028","Suresh Kumar"],
        ["M-110","Grease bearings","3 months","01-Jun-2026","01-Sep-2026","Suresh Kumar"],
        ["L-401","Pigging","Quarterly","22-Apr-2026","22-Jul-2026","Meena Iyer"],
        ["L-401","Cathodic protection survey","Annual","10-Mar-2026","10-Mar-2027","Meena Iyer"],
    ]
    for row_data in data:
        ws.append(row_data)
    for row in ws.iter_rows(min_row=2, max_row=ws.max_row, max_col=len(headers)):
        for cell in row:
            cell.border = thin_border
    ws.column_dimensions['A'].width = 12; ws.column_dimensions['B'].width = 35
    ws.column_dimensions['C'].width = 14; ws.column_dimensions['D'].width = 14
    ws.column_dimensions['E'].width = 14; ws.column_dimensions['F'].width = 18
    wb.save(make_output("Excel", "Maintenance_Schedule.xlsx"))
    files.append("Excel/Maintenance_Schedule.xlsx")

    # --- Spare Parts Inventory ---
    wb = Workbook(); ws = wb.active; ws.title = "Spare Parts"
    headers = ["Part No.","Description","For Equipment","Qty in Stock","Min Stock","Unit","Last Ordered"]
    ws.append(headers); style_header(ws)
    data = [
        ["SEAL-5610-45","Mechanical Seal (John Crane 5610, 45mm)","P-102",2,1,"pcs","15-Mar-2026"],
        ["BRG-6319","Bearing SKF 6319 C3","P-101",3,2,"pcs","20-Apr-2026"],
        ["BRG-6320","Bearing SKF 6320 C3","P-101",3,2,"pcs","20-Apr-2026"],
        ["GR-MOB-PEM","Mobil Polyrex EM Grease (400g)","P-101 / M-110",8,4,"cartridges","10-May-2026"],
        ["FLT-100","Air Filter Element","C-201",6,3,"pcs","05-Dec-2025"],
        ["OIL-RX-ND","Atlas Copco Oil Roto-Xtend ND (5L)","C-201",4,2,"pails","20-Jan-2026"],
        ["GKT-B101","Boiler Door Gasket Set","B-101",1,1,"set","10-Feb-2026"],
        ["DIA-220-EPDM","Actuator Diaphragm (EPDM)","V-220",2,1,"pcs","15-Jul-2026"],
        ["TUBE-E301","Heat Exchanger Tube (SS316L, 19mm OD)","E-301",50,20,"pcs","30-Aug-2024"],
        ["BELT-CT01","V-Belt Set (SPB-2500)","CT-01",12,6,"pcs","05-Apr-2026"],
        ["SEAL-KIT-PIP","Pipeline Gasket Kit 12\"","L-401",3,2,"set","10-Mar-2026"],
        ["MTR-110-FAN","Fan Impeller (M-110 spare)","M-110",1,1,"pcs","01-Sep-2023"],
    ]
    for row_data in data:
        ws.append(row_data)
    for row in ws.iter_rows(min_row=2, max_row=ws.max_row, max_col=len(headers)):
        for cell in row:
            cell.border = thin_border
    ws.column_dimensions['A'].width = 18; ws.column_dimensions['B'].width = 45
    ws.column_dimensions['C'].width = 18; ws.column_dimensions['D'].width = 14
    ws.column_dimensions['E'].width = 10; ws.column_dimensions['F'].width = 12
    ws.column_dimensions['G'].width = 14
    wb.save(make_output("Excel", "Spare_Parts_Inventory.xlsx"))
    files.append("Excel/Spare_Parts_Inventory.xlsx")
    return files

# ══════════════════════════════════════════════════════════════════
#  8. PowerPoint Presentations (2 PPTX)
# ══════════════════════════════════════════════════════════════════

def generate_powerpoint():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    files = []

    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

    def add_slide(prs, title_text, body_lines, bg_color=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        if bg_color:
            bg = slide.background; fill = bg.fill; fill.solid()
            fill.fore_color.rgb = bg_color
        from pptx.util import Inches, Pt
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1))
        tf = txBox.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = title_text; p.font.size = Pt(36); p.font.bold = True
        p.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))
        tf2 = txBox2.text_frame; tf2.word_wrap = True
        for i, line in enumerate(body_lines):
            if i == 0:
                p2 = tf2.paragraphs[0]
            else:
                p2 = tf2.add_paragraph()
            p2.text = line; p2.font.size = Pt(20); p2.space_after = Pt(8)
        return slide

    # PPT-001 Monthly Safety Training
    add_slide(prs, "Monthly Safety Training  --  July 2026", [
        "ABC Petrochemicals Pvt. Ltd. | Gujarat Manufacturing Plant",
        "",
        f"Presented By: Deepak Joshi (Safety Officer)",
        f"Date: 10-Jul-2026",
    ])
    add_slide(prs, "Agenda", [
        "1. Review of Last Month's Safety Statistics",
        "2. Incident Spotlight: P-101 Oil Leak (INC-001)",
        "3. Near-Miss Reporting Updates",
        "4. Confined Space Entry Refresher",
        "5. PPE Compliance Observations",
        "6. Open Forum / Q&A",
    ])
    add_slide(prs, "Safety Statistics  --  YTD 2026", [
        "Total Recordable Incidents: 2",
        "   • P-101 Oil Leak (Apr 2026)  --  environmental",
        "   • V-220 Valve Failure (Feb 2026)  --  process safety",
        "Lost Time Injuries: 0",
        "Near-Misses Reported: 14 (target: 20/year)",
        "",
        "Current LTI-free days: 527",
        "Target: 1,000 LTI-free days by Dec 2027",
    ])
    add_slide(prs, "Incident Spotlight: P-101 Oil Leak", [
        "Date: 03-Apr-2026 | Equipment: P-101",
        "~2L oil leaked from bearing housing drain plug",
        "Root Cause: Plug not tightened after bearing replacement",
        "Corrective Actions Implemented:",
        "   • Torque verification added to checklist",
        "   • Lock wire applied to drain plug",
        "",
        "Key Learning: Always re-check maintenance actions before sign-off.",
    ])
    add_slide(prs, "Confined Space Entry Refresher", [
        "Permit Types: Hot Work / Cold Work / Confined Space",
        "Gas Testing Required Before Entry:",
        "   • O₂: 19.5%  -  23.5%",
        "   • LEL: < 10%",
        "   • H₂S: < 10 ppm",
        "   • CO: < 25 ppm",
        "",
        "Standby Person Must Be Present At All Times.",
        "Emergency Rescue Plan Must Be Reviewed Before Entry.",
    ])
    add_slide(prs, "PPE Compliance  --  July 2026", [
        "Department-wise Compliance (spot checks):",
        "   Operations: 96%",
        "   Maintenance: 92%",
        "   Contractors: 88%",
        "",
        "Action: Contractor PPE awareness session scheduled.",
        "Target: All departments >= 95% by August 2026.",
    ])
    add_slide(prs, "Thank You", [
        "Report hazards. Look out for each other.",
        "Safety is everyone's responsibility.",
        "",
        "Next Training: August 2026  --  Emergency Evacuation Drill",
    ])
    prs.save(make_output("PowerPoint", "PPT-001_Monthly_Safety_Training_Jul2026.pptx"))
    files.append("PowerPoint/PPT-001_Monthly_Safety_Training_Jul2026.pptx")

    # PPT-002 Plant Overview
    prs2 = Presentation()
    prs2.slide_width = Inches(13.333); prs2.slide_height = Inches(7.5)
    add_slide(prs2, "ABC Petrochemicals Pvt. Ltd.", [
        "Gujarat Manufacturing Plant  --  Plant Overview",
        "",
        "Presented By: Rajesh Mehta (Plant Director)",
        f"Date: 01-Jun-2026",
    ])
    add_slide(prs2, "Plant Profile", [
        "Location: Gujarat, India",
        "Established: 1998 (Major Expansion: 2015)",
        "Area: 85 hectares",
        "Employees: 450 (incl. 120 contract staff)",
        "",
        "Products: Ethylene, Propylene, Polyethylene Intermediates",
        "Annual Capacity: 250,000 MT ethylene",
    ])
    add_slide(prs2, "Key Plant Units", [
        "1. Ethylene Cracker Unit (F-101 furnace)",
        "2. Fractionation Unit (C-4101 deethanizer)",
        "3. Utilities (B-101 boiler, CT-01 cooling tower, C-201 compressor)",
        "4. Tank Farm (T-501, T-3301 storage)",
        "5. Pipeline Network (L-401)",
        "",
        "Support Infrastructure:",
        "   • DCS: Yokogawa CENTUM VP",
        "   • SIS: Triconex (SIL-2 / SIL-3)",
        "   • Maintenance Workshop & Warehouse",
    ])
    add_slide(prs2, "Key Initiatives  --  2026", [
        "1. Digital Transformation: TRACE platform deployment (in progress)",
        "2. Energy Efficiency: Cracker furnace tube replacement (Q4 2026)",
        "3. Safety: Target 1,000 LTI-free days",
        "4. Maintenance: Predictive analytics pilot on C-201 compressor",
        "5. Environmental: Zero liquid discharge (ZLD) feasibility study",
    ])
    add_slide(prs2, "Acknowledgements", [
        "Operations Team: Vikram Singh, Sanjay Patel",
        "Engineering: Priya Sharma",
        "Maintenance: Suresh Kumar",
        "HSE: Deepak Joshi, Ananya Gupta",
        "Quality: Meena Iyer",
        "IT & Systems: Amit Patel",
    ])
    prs2.save(make_output("PowerPoint", "PPT-002_Plant_Overview.pptx"))
    files.append("PowerPoint/PPT-002_Plant_Overview.pptx")
    return files

# ══════════════════════════════════════════════════════════════════
#  9. Scanned PDFs (2)
# ══════════════════════════════════════════════════════════════════

def generate_scanned_pdfs():
    """Generate PDFs that appear as scanned/inspection documents."""
    from fpdf import FPDF
    files = []

    # Helper: "scanned" look
    class ScannedPDF(FPDF):
        def header(self):
            self.set_font("Courier", "", 10)
            self.cell(0, 8, "ABC Petrochemicals Pvt. Ltd. - Gujarat Manufacturing Plant", align="C", ln=True)

    # SCN-001 Inspection Checklist
    pdf = ScannedPDF()
    pdf.add_page()
    pdf.set_font("Courier", "", 12)
    pdf.cell(0, 10, "SCN-001: SAFETY INSPECTION CHECKLIST", ln=True, align="C")
    pdf.ln(5)
    pdf.set_font("Courier", "", 10)
    lines = [
        "Area: Cooling Tower CT-01",
        "Inspector: Deepak Joshi",
        "Date: 05-Apr-2026",
        "Time: 10:30 hrs",
        "",
        "Item                                          OK  NOK  N/A",
        "---------------------------------------------------------",
        "1. Guardrails and handrails                    [Y]  [ ]  [ ]",
        "2. Fan guard condition                         [Y]  [ ]  [ ]",
        "3. Water distribution nozzles clear            [Y]  [ ]  [ ]",
        "4. Fill media condition                        [N]  [Y]  [ ]",
        "     >> Moderate scaling in Cell #3",
        "5. Basin water level                           [Y]  [ ]  [ ]",
        "6. Gearbox oil level                           [Y]  [ ]  [ ]",
        "7. V-belt tension                              [Y]  [ ]  [ ]",
        "8. Electrical panel seal                       [Y]  [ ]  [ ]",
        "9. Emergency stop button functional            [Y]  [ ]  [ ]",
        "10. LOTO point identified                      [Y]  [ ]  [ ]",
        "",
        "---------------------------------------------------------",
        "TOTAL OK:  9    |    NOT OK:  1",
        "",
        "",
        "Inspector Signature: ___Deepak Joshi___",
        "Reviewed By: ________Rajesh Mehta_________",
        "",
        "--- End of Checklist ---",
    ]
    for line in lines:
        pdf.cell(0, 6, line, ln=True)
    pdf.output(make_output("Scanned", "SCN-001_Safety_Inspection_Checklist.pdf"))
    files.append("Scanned/SCN-001_Safety_Inspection_Checklist.pdf")

    # SCN-002 P&ID Drawing (text representation)
    pdf2 = ScannedPDF()
    pdf2.add_page()
    pdf2.set_font("Courier", "", 10)
    pdf2.cell(0, 10, "SCN-002: P&ID - COOLING WATER SUPPLY SYSTEM", ln=True, align="C")
    pdf2.cell(0, 8, "Sheet 1 of 1 | Drawing No. PID-CW-001 | Rev B", ln=True, align="C")
    pdf2.ln(5)
    lines2 = [
        "LEGEND:",
        "  --|-->  Gate Valve      ||  Pump         PSV  Safety Valve",
        "  --|<--  Check Valve     |||  Motor        PIC  Pressure Controller",
        "  ==|==  Control Valve    HX  Heat Exchanger  FT   Flow Transmitter",
        "",
        "",
        "                          COOLING TOWER CT-01",
        "                         +------------------+",
        "              V-101      |  6 cells / 4 fans |",
        "  FROM PROCESS ---|>----+|  Basin            |",
        "                     |   |  T = 31 C         |",
        "                     |   +--------+---------+",
        "                     |            |",
        "                     |            | P-101",
        "                     |            v",
        "                     |     +------------+",
        "                     |     |  [PSV-101] |",
        "                     |     |            |",
        "                     +---->|  P-101     |",
        "                           |  Pump      |",
        "                           +-----+------+",
        "                                 |",
        "                           FT-101 |",
        "                                 v",
        "                          +--------------+",
        "                          |   E-301      |",
        "                          |  Heat        |",
        "                          |  Exchanger   |",
        "                          +------+-------+",
        "                                 |",
        "                          PIC-101 |",
        "                          ==|==   | V-220",
        "                                 v",
        "                           TO PROCESS",
        "",
        "",
        "NOTES:",
        "1. All carbon steel piping, Schedule 40.",
        "2. Insulation required on lines > 50 C.",
        "3.. PSV-101 set at 6.0 barg (pump discharge).",
        "",
        "Drawn By: Priya Sharma     |    Date: 10-Jan-2020",
        "Checked By: Rajesh Mehta  |    Last Rev: 15-Mar-2025",
        "",
        "UNCONTROLLED WHEN PRINTED",
    ]
    for line in lines2:
        pdf2.cell(0, 5, line, ln=True)
    pdf2.output(make_output("Scanned", "SCN-002_P&ID_Cooling_Water.pdf"))
    files.append("Scanned/SCN-002_P&ID_Cooling_Water.pdf")
    return files

# ══════════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    print("=" * 60)
    print("  ABC Petrochemicals  --  Demo Dataset Generator")
    print("=" * 60)

    all_files = []
    sop_files = generate_sops(); all_files.extend(sop_files); print(f"  [OK] SOPs ({len(sop_files)} files)")
    man_files = generate_manuals(); all_files.extend(man_files); print(f"  [OK] Manuals ({len(man_files)} files)")
    mnt_files = generate_maintenance(); all_files.extend(mnt_files); print(f"  [OK] Maintenance ({len(mnt_files)} files)")
    ins_files = generate_inspections(); all_files.extend(ins_files); print(f"  [OK] Inspections ({len(ins_files)} files)")
    inc_files = generate_incidents(); all_files.extend(inc_files); print(f"  [OK] Incidents ({len(inc_files)} files)")
    log_files = generate_shiftlogs(); all_files.extend(log_files); print(f"  [OK] Shift Logs ({len(log_files)} files)")
    xl_files = generate_excel(); all_files.extend(xl_files); print(f"  [OK] Excel Files ({len(xl_files)} files)")
    ppt_files = generate_powerpoint(); all_files.extend(ppt_files); print(f"  [OK] PowerPoint ({len(ppt_files)} files)")
    scn_files = generate_scanned_pdfs(); all_files.extend(scn_files); print(f"  [OK] Scanned PDFs ({len(scn_files)} files)")

    print()
    print("=" * 60)
    print(f"  GENERATION COMPLETE: {len(all_files)} files created")
    print("=" * 60)
    print()
    for f in sorted(all_files):
        print(f"    {f}")
